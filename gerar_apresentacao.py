from pptx import Presentation
from pptx.util import Inches
import pandas as pd
import os

# ==========================================
# CONFIGURAÇÕES
# ==========================================

os.makedirs("relatorio", exist_ok=True)

PPT_SAIDA = "relatorio/apresentacao_executiva.pptx"

# ==========================================
# LEITURA DOS DADOS
# ==========================================

try:
    df = pd.read_csv("data/vendas.csv")

    faturamento_total = round(
        df["valor_total"].sum(),
        2
    )

    total_vendas = len(df)

    ticket_medio = round(
        faturamento_total / total_vendas,
        2
    )

    satisfacao = round(
        df["avaliacao"].mean(),
        2
    )

except Exception:

    faturamento_total = 0
    total_vendas = 0
    ticket_medio = 0
    satisfacao = 0

# ==========================================
# APRESENTAÇÃO
# ==========================================

prs = Presentation()

# ==========================================
# FUNÇÃO AUXILIAR
# ==========================================

def adicionar_slide(titulo, conteudo):

    slide = prs.slides.add_slide(
        prs.slide_layouts[1]
    )

    slide.shapes.title.text = titulo

    slide.placeholders[1].text = conteudo

# ==========================================
# SLIDE 1
# ==========================================

slide = prs.slides.add_slide(
    prs.slide_layouts[0]
)

slide.shapes.title.text = (
    "Tressenza Big Data Analytics"
)

slide.placeholders[1].text = (
    "Relatório Executivo e Estratégico"
)

# ==========================================
# SLIDE 2
# ==========================================

adicionar_slide(
    "Sobre a Empresa",
    """
Empresa de médio porte do setor
de perfumaria e cosméticos.

Atuação física e digital.
Foco em crescimento sustentável.
"""
)

# ==========================================
# SLIDE 3
# ==========================================

adicionar_slide(
    "Visão Geral das Operações",
    f"""
Total de vendas:
{total_vendas:,}

Faturamento total:
R$ {faturamento_total:,.2f}
"""
)

# ==========================================
# SLIDE 4
# ==========================================

adicionar_slide(
    "Indicadores Financeiros",
    f"""
Ticket Médio:
R$ {ticket_medio:,.2f}

Receita crescente
ao longo dos períodos analisados.
"""
)

# ==========================================
# SLIDE 5
# ==========================================

adicionar_slide(
    "Categorias Mais Vendidas",
    """
Perfumes

Maquiagem

Skincare

Cabelos

Kits Presente
"""
)

# ==========================================
# SLIDE 6
# ==========================================

adicionar_slide(
    "Produtos Destaque",
    """
Perfume Luxury

Perfume Essence

Creme Facial Premium

Kit Presente Gold
"""
)

# ==========================================
# SLIDE 7
# ==========================================

adicionar_slide(
    "Satisfação dos Clientes",
    f"""
Avaliação Média:

{satisfacao}/5

Alta taxa de aprovação.
"""
)

# ==========================================
# SLIDE 8
# ==========================================

adicionar_slide(
    "Estratégias de Marketing",
    """
Instagram

TikTok

Programa de Fidelidade

Cupons de Desconto

Campanhas Sazonais
"""
)

# ==========================================
# SLIDE 9
# ==========================================

adicionar_slide(
    "Machine Learning",
    """
Aplicação de Regressão Linear.

Análise de tendências.

Previsão de crescimento
para os próximos 5 anos.
"""
)

# ==========================================
# SLIDE 10
# ==========================================

adicionar_slide(
    "Projeção para 5 Anos",
    """
Expansão do faturamento.

Ampliação da base de clientes.

Novos canais de venda.

Maior participação digital.
"""
)

# ==========================================
# SLIDE 11
# ==========================================

adicionar_slide(
    "Plano de Expansão",
    """
Novas unidades físicas.

Marketplace próprio.

Programa VIP.

Expansão nacional.
"""
)

# ==========================================
# SLIDE 12
# ==========================================

adicionar_slide(
    "Conclusão",
    """
O negócio apresenta
crescimento consistente.

Os dados indicam potencial
de expansão sustentável
nos próximos anos.
"""
)

# ==========================================
# INSERÇÃO DE GRÁFICOS
# ==========================================

graficos = [
    "graficos/faturamento_mensal.png",
    "graficos/categorias.png",
    "graficos/previsao_5_anos.png"
]

for grafico in graficos:

    if os.path.exists(grafico):

        slide = prs.slides.add_slide(
            prs.slide_layouts[5]
        )

        slide.shapes.title.text = (
            os.path.basename(grafico)
        )

        slide.shapes.add_picture(
            grafico,
            Inches(0.5),
            Inches(1),
            width=Inches(8)
        )

# ==========================================
# SALVAR
# ==========================================

prs.save(PPT_SAIDA)

print("=" * 60)
print("APRESENTAÇÃO GERADA")
print("=" * 60)
print(PPT_SAIDA)
print("=" * 60)
